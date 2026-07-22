"""
sales_growth/build_ppt.py
Built-out consolidated deck: title, a 3-tier overview (stacked bands so all 13
initiatives fit), one detail slide per initiative (key decisions + 90-day + year-one,
adaptive font so text never overflows), and a closing / next-steps slide.
Content comes from plan_data.PLAN so the deck, the Excel, and the List seed agree.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from plan_data import PLAN, rollup, TIER_ORDER, MA_TARGETS

CORE  = RGBColor(0x00,0x67,0xB1); EXPAND=RGBColor(0xB8,0x79,0x1F); ENABLE=RGBColor(0x2E,0x7B,0x87)
INK   = RGBColor(0x1F,0x3A,0x4D); MUTE  = RGBColor(0x5A,0x6B,0x7A); MUTE2=RGBColor(0x6B,0x7B,0x8A)
LIGHT = RGBColor(0xF1,0xF5,0xF9); WHITE=RGBColor(0xFF,0xFF,0xFF); PANEL=RGBColor(0xF7,0xF9,0xFB)
LIME  = RGBColor(0x99,0xC2,0x21)
FONT  = "Segoe UI"
LOGO  = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "static", "ieg-logo.png")
TIER_NAME = {"Core Play":"CORE PLAYS","Expansion Bet":"EXPANSION BETS","Enabler":"ENABLERS"}
TIER_COLOR={"Core Play":CORE,"Expansion Bet":EXPAND,"Enabler":ENABLE}
TIER_LABEL={"Core Play":"TIER A  ·  CORE PLAY","Expansion Bet":"TIER B  ·  EXPANSION BET","Enabler":"TIER C  ·  ENABLER"}
SW, SH = Inches(13.333), Inches(7.5)


def _box(sl,l,t,w,h,fill=None,line=None):
    sh=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h); sh.shadow.inherit=False
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(0.75)
    return sh


def _txt(sl,l,t,w,h,runs,size=12,bold=False,color=INK,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,
         space=4,line_spacing=1.0):
    tb=sl.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Pt(2); tf.margin_right=Pt(2); tf.margin_top=Pt(1); tf.margin_bottom=Pt(1)
    if isinstance(runs,str): runs=[runs]
    for k,item in enumerate(runs):
        p=tf.paragraphs[0] if k==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space); p.space_before=Pt(0); p.line_spacing=line_spacing
        text,opts=(item if isinstance(item,tuple) else (item,{}))
        r=p.add_run(); r.text=text
        r.font.size=Pt(opts.get("size",size)); r.font.bold=opts.get("bold",bold)
        r.font.name=FONT; r.font.color.rgb=opts.get("color",color)
    return tb


def _inline(sl,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    """One paragraph, multiple runs (for meta / KPI lines that must stay on one line)."""
    tb=sl.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Pt(2); tf.margin_right=Pt(2); tf.margin_top=Pt(1); tf.margin_bottom=Pt(1)
    p=tf.paragraphs[0]; p.alignment=align; p.line_spacing=1.0
    for text,opts in runs:
        r=p.add_run(); r.text=text
        r.font.size=Pt(opts.get("size",11)); r.font.bold=opts.get("bold",False)
        r.font.name=FONT; r.font.color.rgb=opts.get("color",INK)
    return tb


def _bullets(sl,l,t,w,h,items,base=12,minsz=9,color=INK,space=5):
    """Render '•' bullets, shrinking the font by count/length so they fit the region."""
    n=len(items); longest=max((len(x) for x in items),default=0)
    size=base
    if n>=5: size=base-2
    elif n==4: size=base-1
    if longest>92 and size>minsz: size-=1
    size=max(size,minsz)
    _txt(sl,l,t,w,h,[(f"•  {x}",{"size":size,"color":color}) for x in items],space=space,line_spacing=1.0)


def _footer(sl,page):
    _box(sl,0,SH-Inches(0.4),SW,Inches(0.4),fill=PANEL)
    _txt(sl,Inches(0.5),SH-Inches(0.37),Inches(10),Inches(0.3),
         "IEG Sales Growth  ·  Consolidated Plan  ·  feeds the SLT Strategic Initiatives tracker",size=8.5,color=MUTE2)
    _txt(sl,SW-Inches(1.1),SH-Inches(0.37),Inches(0.8),Inches(0.3),str(page),size=8.5,color=MUTE2,align=PP_ALIGN.RIGHT)


def _blank(prs):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); _box(sl,0,0,SW,SH,fill=WHITE); return sl


def _money(v):
    return f"${v/1e6:.1f}M" if v else None


# ---------------------------------------------------------------- title
def title_slide(prs):
    sl=_blank(prs); r=rollup()
    try: sl.shapes.add_picture(LOGO, Inches(0.85), Inches(0.55), height=Inches(0.6))
    except Exception: pass
    _txt(sl,Inches(0.9),Inches(1.5),Inches(11),Inches(0.4),
         "IEG SALES GROWTH STRATEGY WORKSHOP  ·  JULY 2026",size=13,bold=True,color=CORE)
    _txt(sl,Inches(0.85),Inches(1.92),Inches(11.6),Inches(1.05),[("The Consolidated Plan",{"size":44,"bold":True,"color":INK})])
    _box(sl,Inches(0.9),Inches(2.98),Inches(1.5),Inches(0.06),fill=LIME)
    _txt(sl,Inches(0.9),Inches(3.2),Inches(11.4),Inches(0.75),
         "Every idea from the six SLT decks, consolidated into 13 initiatives across three tiers, with the decisions "
         "and the 90-day and year-one actions to move each one.",size=15,color=MUTE)
    cards=[("Core Play",CORE,"the proven consensus"),("Expansion Bet",EXPAND,"bigger bets, to size"),("Enabler",ENABLE,"margin + enabling")]
    x0=Inches(0.9); cw=Inches(3.7); gap=Inches(0.38); top=Inches(4.3); ch=Inches(1.9)
    for k,(tier,col,tag) in enumerate(cards):
        x=Emu(int(x0)+k*(int(cw)+int(gap)))
        _box(sl,x,top,cw,ch,fill=LIGHT); _box(sl,x,top,cw,Inches(0.48),fill=col)
        _txt(sl,x,top,cw,Inches(0.48),TIER_NAME[tier],size=13,bold=True,color=WHITE,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        amt=r["by_tier"][tier]; big=f"${amt/1e6:.1f}M" if amt else "TBD"
        _txt(sl,x,top+Inches(0.62),cw,Inches(0.72),[(big,{"size":38,"bold":True,"color":col})],align=PP_ALIGN.CENTER)
        _txt(sl,x,top+Inches(1.42),cw,Inches(0.4),f"{r['counts'][tier]} initiatives · {tag}",size=11,color=MUTE2,align=PP_ALIGN.CENTER)
    _box(sl,Inches(0.9),Inches(6.5),Inches(11.53),Inches(0.55),fill=CORE)
    _inline(sl,Inches(1.15),Inches(6.5),Inches(11.1),Inches(0.55),
        [(f"~${r['total']/1e6:.0f}M annual revenue target across the growth plays",{"size":14,"bold":True,"color":WHITE}),
         ("      M&A sized in diligence  ·  feeds the SLT Strategic Initiatives tracker",{"size":11,"color":RGBColor(0xCF,0xE3,0xF3)})],anchor=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------- overview
def _ov_card(sl,x,y,w,h,i,col,compact=False):
    _box(sl,x,y,w,h,fill=LIGHT); _box(sl,x,y,Inches(0.07),h,fill=col)
    money=_money(i["rev"]) or "TBD"
    nm=i["name"]
    _txt(sl,x+Inches(0.16),y+Inches(0.05),w-Inches(0.24),Inches(0.55) if compact else Inches(0.5),
         [(f"{i['code']}  ",{"size":10 if compact else 11,"bold":True,"color":col}),
          (nm,{"size":8.5 if compact else 10.5,"bold":True,"color":INK})],line_spacing=0.92)
    _txt(sl,x+Inches(0.16),y+h-Inches(0.32),w-Inches(0.24),Inches(0.3),
         [(f"S {i['sponsor'].split()[0]} · O {i['owner'].split()[0]}   ",{"size":8,"color":MUTE2}),
          (money,{"size":8.5,"bold":True,"color":col})],line_spacing=0.9)


def overview_slide(prs,page):
    sl=_blank(prs)
    _txt(sl,Inches(0.5),Inches(0.3),Inches(12),Inches(0.55),[("The plan at a glance",{"size":30,"bold":True,"color":INK})])
    _txt(sl,Inches(0.5),Inches(0.92),Inches(12.3),Inches(0.5),
         "All six leaders independently prioritized selling deeper into existing customers and growing recurring service. "
         "Core plays first (the consensus), expansion bets second, enablers underneath.",size=11.5,color=MUTE)
    rr=rollup()
    def band(tier,ytop,cols,ch):
        col=TIER_COLOR[tier]; items=[z for z in PLAN if z["tier"]==tier]
        amt=rr["by_tier"][tier]; sub=(f"   ${amt/1e6:.1f}M roll-up" if amt else "   TBD")
        _inline(sl,Inches(0.5),ytop,Inches(12),Inches(0.28),
                [(TIER_LABEL[tier],{"size":11,"bold":True,"color":col}),(sub,{"size":11,"bold":True,"color":MUTE2})])
        usable=int(SW)-int(Inches(1.0)); gap=int(Inches(0.14))
        cw=(usable-gap*(cols-1))//cols; x0=int(Inches(0.5)); yy=int(ytop)+int(Inches(0.34))
        for k,i in enumerate(items):
            x=Emu(x0+k*(cw+gap)); _ov_card(sl,x,Emu(yy),Emu(cw),ch,i,col,compact=(cols>=5))
    band("Core Play",Inches(1.5),3,Inches(0.9))
    band("Expansion Bet",Inches(2.94),6,Inches(1.0))
    band("Enabler",Inches(4.46),4,Inches(0.9))
    # consensus strip
    _box(sl,Inches(0.5),Inches(5.85),Inches(12.33),Inches(1.1),fill=PANEL)
    _txt(sl,Inches(0.7),Inches(5.95),Inches(12),Inches(0.4),[("Where the whole room agreed",{"size":14,"bold":True,"color":INK})])
    _txt(sl,Inches(0.7),Inches(6.33),Inches(12),Inches(0.6),
         "Independently, every leader prioritized Own the Account (attach / cross-sell) and the Service PM channel. The three enablers - "
         "comp redesign, the pricing revamp, and the demand engine - fund and unblock the rest, and start now. M&A is the accelerant, sized in diligence.",
         size=12,color=MUTE,line_spacing=1.1)
    _footer(sl,page)


# ---------------------------------------------------------------- detail
def detail_slide(prs,i,page):
    sl=_blank(prs); col=TIER_COLOR[i["tier"]]
    _box(sl,0,0,SW,Inches(1.28),fill=col)
    _txt(sl,Inches(0.5),Inches(0.16),Inches(8),Inches(0.28),TIER_LABEL[i["tier"]],size=11,bold=True,color=WHITE)
    tname=f"{i['code']}   {i['name']}"
    _txt(sl,Inches(0.5),Inches(0.44),Inches(8.5),Inches(0.82),[(tname,{"size":22,"bold":True,"color":WHITE})],line_spacing=0.95)
    money=_money(i["rev"]); chip=(money+" revenue target") if money else "Target sized after 90-day analysis"
    if i["ebitda"]: chip+=f"  ·  {_money(i['ebitda'])} GM/EBITDA"
    _txt(sl,SW-Inches(4.25),Inches(0.2),Inches(3.85),Inches(0.85),chip,size=12,bold=True,color=WHITE,align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.TOP)

    _inline(sl,Inches(0.5),Inches(1.4),Inches(12.3),Inches(0.32),
         [("Sponsor: ",{"size":11,"color":MUTE2}),(i["sponsor"],{"size":11,"bold":True,"color":INK}),
          ("      Owner: ",{"size":11,"color":MUTE2}),(i["owner"],{"size":11,"bold":True,"color":INK}),
          ("      Priority: ",{"size":11,"color":MUTE2}),(i["priority"],{"size":11,"bold":True,"color":INK}),
          ("      Region: ",{"size":11,"color":MUTE2}),(i["region"],{"size":11,"bold":True,"color":INK}),
          ("      Contributors: ",{"size":11,"color":MUTE2}),(", ".join(n.split()[0] for n in i["contributors"]),{"size":11,"color":INK})])
    _txt(sl,Inches(0.5),Inches(1.78),Inches(12.3),Inches(0.78),i["desc"],size=11,color=MUTE,line_spacing=1.05)

    hy=Inches(2.62)
    # left: decisions
    _box(sl,Inches(0.5),hy,Inches(4.35),Inches(0.32),fill=col)
    _txt(sl,Inches(0.5),hy,Inches(4.35),Inches(0.32),"KEY DECISIONS",size=11.5,bold=True,color=WHITE,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    _bullets(sl,Inches(0.55),hy+Inches(0.45),Inches(4.3),Inches(3.4),i["decisions"],base=12,minsz=10,space=9)
    # right: 90-day
    rx=Inches(5.2)
    _box(sl,rx,hy,Inches(7.6),Inches(0.32),fill=INK)
    _txt(sl,rx,hy,Inches(7.6),Inches(0.32),"FIRST 90 DAYS",size=11.5,bold=True,color=WHITE,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    d90=[n for n,p,_ in i["tasks"] if p=="90d"]; dyr=[n for n,p,_ in i["tasks"] if p=="year"]
    _bullets(sl,rx+Inches(0.05),hy+Inches(0.45),Inches(7.55),Inches(1.85),d90,base=11.5,minsz=9,space=6)
    yy=Inches(5.02)
    _box(sl,rx,yy,Inches(7.6),Inches(0.32),fill=MUTE2)
    _txt(sl,rx,yy,Inches(7.6),Inches(0.32),"YEAR ONE",size=11.5,bold=True,color=WHITE,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    _bullets(sl,rx+Inches(0.05),yy+Inches(0.45),Inches(7.55),Inches(1.0),dyr,base=11.5,minsz=9,space=6)
    # KPI
    _box(sl,Inches(0.5),Inches(6.5),Inches(12.33),Inches(0.48),fill=LIGHT)
    _inline(sl,Inches(0.65),Inches(6.5),Inches(12.05),Inches(0.48),
         [("KPI   ",{"size":11,"bold":True,"color":col}),(i["kpi"],{"size":10.5,"color":INK})],anchor=MSO_ANCHOR.MIDDLE)
    _footer(sl,page)


# ---------------------------------------------------------------- M&A grid
def _cell(cell,text,white=False,bold=False,size=11,color=INK,fill=None):
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb=fill
    cell.margin_left=Pt(9); cell.margin_right=Pt(9); cell.margin_top=Pt(5); cell.margin_bottom=Pt(5)
    cell.vertical_anchor=MSO_ANCHOR.MIDDLE
    tf=cell.text_frame; tf.word_wrap=True; tf.clear()
    p=tf.paragraphs[0]; p.line_spacing=1.0
    r=p.add_run(); r.text=text
    r.font.size=Pt(size); r.font.bold=bold; r.font.name=FONT
    r.font.color.rgb=(WHITE if white else color)


def ma_targets_slide(prs,page):
    sl=_blank(prs); col=ENABLE
    _box(sl,0,0,SW,Inches(1.28),fill=col)
    _txt(sl,Inches(0.5),Inches(0.16),Inches(8),Inches(0.28),"TIER C  ·  ENABLER  ·  ACCELERATE",size=11,bold=True,color=WHITE)
    _txt(sl,Inches(0.5),Inches(0.44),Inches(11),Inches(0.82),[("C12   M&A Tier-A Target Slate",{"size":22,"bold":True,"color":WHITE})],line_spacing=0.95)
    _txt(sl,Inches(0.5),Inches(1.5),Inches(12.3),Inches(0.55),
         "The on-strategy Tier-A slate that makes the M&A block credible. Soft-vetted (desktop diligence only); "
         "revenue, EBITDA and EV are deliberately left off here and get sized in diligence.",size=12,color=MUTE,line_spacing=1.05)
    rows=len(MA_TARGETS)+1
    tbl=sl.shapes.add_table(rows,3,Inches(0.5),Inches(2.35),Inches(12.33),Inches(3.5)).table
    tbl.first_row=False; tbl.horz_banding=False
    tbl.columns[0].width=Inches(2.35); tbl.columns[1].width=Inches(2.45); tbl.columns[2].width=Inches(7.53)
    for c,h in enumerate(["Target","Serves","Strategic fit"]):
        _cell(tbl.cell(0,c),h,white=True,bold=True,size=12,fill=INK)
    tbl.rows[0].height=Inches(0.4)
    for r,(name,serves,fit) in enumerate(MA_TARGETS,1):
        _cell(tbl.cell(r,0),name,bold=True,size=13,color=col,fill=(LIGHT if r%2 else WHITE))
        _cell(tbl.cell(r,1),serves,size=11.5,color=INK,fill=(LIGHT if r%2 else WHITE))
        _cell(tbl.cell(r,2),fit,size=11.5,color=INK,fill=(LIGHT if r%2 else WHITE))
        tbl.rows[r].height=Inches(0.72)
    _txt(sl,Inches(0.5),Inches(6.2),Inches(12.3),Inches(0.5),
         "We would not buy all of them - this is what makes the M&A block credible, not a shopping list. "
         "The slate plugs into A2 (service), B5 (chutes) and B9 (reman / absorbents).",size=11,color=MUTE2,line_spacing=1.05)
    _footer(sl,page)


# ---------------------------------------------------------------- closing
def closing_slide(prs,page):
    sl=_blank(prs); _box(sl,0,0,Inches(0.28),SH,fill=CORE)
    n_task=sum(len(i["tasks"]) for i in PLAN)
    _txt(sl,Inches(0.85),Inches(0.55),Inches(11),Inches(0.7),[("From plan to tracker",{"size":30,"bold":True,"color":INK})])
    steps=[
      ("1.  Load the tracker",f"The 13 initiatives and their {n_task} tasks populate the SLT Strategic Initiatives List (Category BOD, Sponsor / Owner / Task), so progress rolls up on the live dashboard."),
      ("2.  Confirm targets",  "Sponsors and Owners are set. A2 (PM channel) carries a revised ~$4M target; the expansion-bet dollars are first-pass and get sized in the 90-day work. Clare-to-Elgin and M&A stay unsized until their analysis / diligence."),
      ("3.  Work the 90-day list","Each initiative's first-90-day actions are the near-term tasks. The enablers (comp, the pricing revamp, the demand engine) unblock the rest and start now. The Account Health + CO worklists are already built and go in first."),
      ("4.  Resolve the tensions","Pricing enforcement vs rebates; online containers vs the OEM channel; the core-parts aperture; discipline vs the bigger platform bets. Decide these to make the plan one voice."),
      ("5.  Report up","The dashboard's board view rolls the initiatives to the BOD; the tracker keeps Sponsor / Owner accountable between meetings."),
    ]
    y=Inches(1.6)
    for head,body in steps:
        _txt(sl,Inches(0.9),y,Inches(3.2),Inches(0.8),[(head,{"size":15,"bold":True,"color":CORE})])
        _txt(sl,Inches(4.2),y,Inches(8.5),Inches(0.95),body,size=12.5,color=INK,line_spacing=1.05)
        y+=Inches(1.04)
    _footer(sl,page)


def build(path):
    prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH
    title_slide(prs)
    overview_slide(prs,2)
    page=3
    for i in PLAN:
        detail_slide(prs,i,page); page+=1
    ma_targets_slide(prs,page); page+=1
    closing_slide(prs,page)
    try: prs.save(path)
    except PermissionError:
        path=path.replace(".pptx","_v2.pptx"); prs.save(path)
    return path,len(prs.slides._sldIdLst)


if __name__ == "__main__":
    out=os.path.join(os.path.expanduser("~"),"Downloads","IEG Sales Growth - Consolidated Plan.pptx")
    p,n=build(out); print(f"saved {p}  ({n} slides)")
